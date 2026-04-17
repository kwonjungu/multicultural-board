"""
Build teacher manual for 🐝 꿀벌 소통창.
Run: python scripts/build-manual-pptx.py
Output: docs/teacher-manual.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "teacher-manual.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# ---------- Theme ----------
HONEY_GOLD = RGBColor(0xF5, 0x9E, 0x0B)
HONEY_LIGHT = RGBColor(0xFB, 0xBF, 0x24)
HONEY_BG = RGBColor(0xFE, 0xF3, 0xC7)
INK = RGBColor(0x11, 0x18, 0x27)
SUB = RGBColor(0x6B, 0x72, 0x80)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFA, 0xFB, 0xFF)
PURPLE = RGBColor(0x5B, 0x57, 0xF5)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xEF, 0x44, 0x44)
LIGHT_BORDER = RGBColor(0xE5, 0xE7, 0xEB)

FONT = "맑은 고딕"

# ---------- Presentation (16:9) ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_bg(slide, color=BG):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.shadow.inherit = False
    return shp


def rect(slide, x, y, w, h, fill=CARD, line=None, radius=None):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    if radius is not None:
        shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, content, size=18, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = content if isinstance(content, list) else [content]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def chip(slide, x, y, label, fill=HONEY_BG, color=HONEY_GOLD, size=11):
    # Rough pill sized by label
    w = Inches(0.2 + 0.13 * max(6, len(label)))
    h = Inches(0.32)
    shp = rect(slide, x, y, w, h, fill=fill, radius=0.5)
    text(slide, x, y, w, h, label, size=size, bold=True, color=color,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return shp, w, h


def numbered_step(slide, x, y, w, n, title, body, step_color=HONEY_GOLD):
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = step_color
    circle.line.fill.background()
    circle.shadow.inherit = False
    text(slide, x, y, Inches(0.45), Inches(0.45), str(n),
         size=15, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Title + body
    tx = x + Inches(0.6)
    tw = w - Inches(0.6)
    text(slide, tx, y - Inches(0.02), tw, Inches(0.35), title,
         size=15, bold=True, color=INK)
    text(slide, tx, y + Inches(0.32), tw, Inches(0.7), body,
         size=12, color=SUB)


def page_header(slide, chapter_no, chapter_title, subtitle=""):
    # Top honey bar
    bar = rect(slide, 0, 0, SW, Inches(0.85),
               fill=HONEY_GOLD, radius=None)
    # Chapter badge
    text(slide, Inches(0.6), Inches(0.15), Inches(2.2), Inches(0.3),
         f"Chapter {chapter_no}", size=11, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF))
    text(slide, Inches(0.6), Inches(0.38), Inches(10), Inches(0.45),
         chapter_title, size=22, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF))
    # Bee emoji right
    text(slide, SW - Inches(1.1), Inches(0.15), Inches(0.9), Inches(0.6),
         "🐝", size=32, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
             subtitle, size=13, color=SUB)


def footer(slide, page_n, total):
    text(slide, Inches(0.6), SH - Inches(0.4), Inches(6), Inches(0.3),
         "🐝 꿀벌 소통창 · 선생님용 매뉴얼", size=9, color=SUB)
    text(slide, SW - Inches(2), SH - Inches(0.4), Inches(1.4), Inches(0.3),
         f"{page_n} / {total}", size=9, color=SUB, align=PP_ALIGN.RIGHT)


# ---------- Visual mock builders (used in place of screenshots) ----------
def mock_phone(slide, cx, cy, w, h, draw_content):
    """Draw a phone frame and call draw_content(slide, ix, iy, iw, ih)
    for the inner screen area."""
    x = cx - w // 2
    y = cy - h // 2
    # outer frame
    frame = rect(slide, x, y, w, h, fill=INK, radius=0.08)
    pad = Inches(0.08)
    inner_x = x + pad
    inner_y = y + pad
    inner_w = w - 2 * pad
    inner_h = h - 2 * pad
    inner = rect(slide, inner_x, inner_y, inner_w, inner_h, fill=BG, radius=0.04)
    draw_content(slide, inner_x, inner_y, inner_w, inner_h)
    return frame


def mock_screen(slide, x, y, w, h, draw_content):
    """Draw a clean browser/screen rectangle with inner content."""
    # Shadow layer
    shadow = rect(slide, x + Inches(0.05), y + Inches(0.06), w, h,
                  fill=RGBColor(0xD1, 0xD5, 0xDB), radius=0.03)
    # Browser top bar
    top = rect(slide, x, y, w, Inches(0.3),
               fill=RGBColor(0xE5, 0xE7, 0xEB), radius=0.03)
    # Traffic lights
    for i, c in enumerate([RGBColor(0xEF, 0x44, 0x44),
                            RGBColor(0xF5, 0x9E, 0x0B),
                            RGBColor(0x22, 0xC5, 0x5E)]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            x + Inches(0.12 + i * 0.18), y + Inches(0.08),
            Inches(0.13), Inches(0.13))
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.fill.background()
        dot.shadow.inherit = False
    # Content area
    inner_y = y + Inches(0.3)
    inner_h = h - Inches(0.3)
    inner = rect(slide, x, inner_y, w, inner_h, fill=CARD)
    draw_content(slide, x, inner_y, w, inner_h)


# ---------- Specific mocks ----------
def mock_cover_landing(slide, x, y, w, h):
    # Hero center
    cx = x + w // 2
    text(slide, x, y + Inches(0.6), w, Inches(0.6), "🐝",
         size=48, align=PP_ALIGN.CENTER)
    text(slide, x, y + Inches(1.4), w, Inches(0.5), "꿀벌 소통창",
         size=26, bold=True, color=INK, align=PP_ALIGN.CENTER)
    text(slide, x, y + Inches(2.0), w, Inches(0.4),
         "전 세계 어디든, 꿀벌처럼 소식을 나눠요",
         size=12, color=SUB, align=PP_ALIGN.CENTER)
    # Buttons
    btn_w = Inches(3.0)
    btn_h = Inches(0.6)
    bx = cx - btn_w // 2
    rect(slide, bx, y + Inches(2.8), btn_w, btn_h, fill=HONEY_GOLD, radius=0.5)
    text(slide, bx, y + Inches(2.8), btn_w, btn_h, "➕  새 방 만들기",
         size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, bx, y + Inches(3.5), btn_w, btn_h, fill=CARD,
         line=HONEY_GOLD, radius=0.5)
    text(slide, bx, y + Inches(3.5), btn_w, btn_h, "🔑  방 코드 입력",
         size=14, bold=True, color=HONEY_GOLD,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def mock_board(slide, x, y, w, h):
    # Top bar
    rect(slide, x, y, w, Inches(0.6), fill=RGBColor(0xFF, 0xFF, 0xFF))
    text(slide, x + Inches(0.2), y + Inches(0.12), Inches(4), Inches(0.35),
         "🐝 우리반 소통창", size=14, bold=True)
    text(slide, x + Inches(4), y + Inches(0.15), Inches(2), Inches(0.3),
         "방 코드: ABC123", size=10, color=SUB)
    # Right-side chips (mock)
    rect(slide, x + w - Inches(2.7), y + Inches(0.14), Inches(0.9),
         Inches(0.32), fill=HONEY_BG, radius=0.5)
    text(slide, x + w - Inches(2.7), y + Inches(0.14), Inches(0.9),
         Inches(0.32), "🎮 게임", size=9, bold=True, color=HONEY_GOLD,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, x + w - Inches(1.7), y + Inches(0.14), Inches(0.7),
         Inches(0.32), fill=PURPLE, radius=0.5)
    text(slide, x + w - Inches(1.7), y + Inches(0.14), Inches(0.7),
         Inches(0.32), "➕ 카드", size=9, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, x + w - Inches(0.95), y + Inches(0.14), Inches(0.7),
         Inches(0.32), fill=CARD, line=HONEY_GOLD, radius=0.5)
    text(slide, x + w - Inches(0.95), y + Inches(0.14), Inches(0.7),
         Inches(0.32), "⬇ 내보내기", size=9, bold=True, color=HONEY_GOLD,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Cards grid
    card_y = y + Inches(0.8)
    card_h = Inches(1.4)
    card_w = Inches(1.7)
    gap = Inches(0.15)
    cols = 5
    start_x = x + Inches(0.3)
    samples = [
        ("안녕하세요!", "🇰🇷", HONEY_BG),
        ("Xin chào!", "🇻🇳", RGBColor(0xDC, 0xFC, 0xE7)),
        ("你好!", "🇨🇳", RGBColor(0xFE, 0xE2, 0xE2)),
        ("Hello!", "🇺🇸", RGBColor(0xDB, 0xEA, 0xFE)),
        ("こんにちは", "🇯🇵", RGBColor(0xFE, 0xF3, 0xC7)),
    ]
    for i in range(cols):
        cx = start_x + i * (card_w + gap)
        label, flag, fc = samples[i]
        rect(slide, cx, card_y, card_w, card_h, fill=fc, radius=0.08)
        text(slide, cx + Inches(0.1), card_y + Inches(0.1), Inches(0.4),
             Inches(0.4), flag, size=18)
        text(slide, cx + Inches(0.1), card_y + Inches(0.55), card_w - Inches(0.2),
             Inches(0.4), label, size=11, bold=True)
        text(slide, cx + Inches(0.1), card_y + Inches(0.9),
             card_w - Inches(0.2), Inches(0.3),
             "학생1 · 방금", size=8, color=SUB)
    # Second row
    card_y2 = card_y + card_h + Inches(0.2)
    samples2 = [
        ("🎙 음성 메모", HONEY_BG),
        ("📷 사진 카드", RGBColor(0xDB, 0xEA, 0xFE)),
        ("활동지 🖼", RGBColor(0xEF, 0xE4, 0xFD)),
    ]
    for i, (label, fc) in enumerate(samples2):
        cx = start_x + i * (card_w + gap)
        rect(slide, cx, card_y2, card_w, card_h, fill=fc, radius=0.08)
        text(slide, cx, card_y2 + Inches(0.45), card_w, Inches(0.5),
             label, size=12, bold=True, align=PP_ALIGN.CENTER)


def mock_card_modal(slide, x, y, w, h):
    # Dim background
    rect(slide, x, y, w, h, fill=RGBColor(0xE5, 0xE7, 0xEB))
    # Modal center
    mw = Inches(4.5)
    mh = Inches(4.5)
    mx = x + (w - mw) // 2
    my = y + (h - mh) // 2
    rect(slide, mx, my, mw, mh, fill=CARD, radius=0.05)
    text(slide, mx + Inches(0.3), my + Inches(0.2), mw - Inches(0.6),
         Inches(0.4), "✏️ 새 카드 만들기", size=15, bold=True)
    # Author
    rect(slide, mx + Inches(0.3), my + Inches(0.75), mw - Inches(0.6),
         Inches(0.45), fill=BG, radius=0.08)
    text(slide, mx + Inches(0.45), my + Inches(0.85), Inches(3),
         Inches(0.25), "이름 입력", size=11, color=SUB)
    # Language chips
    chipset_y = my + Inches(1.35)
    for i, (flag, lang) in enumerate([("🇰🇷", "한국어"), ("🇻🇳", "Tiếng Việt"),
                                       ("🇨🇳", "中文"), ("🇺🇸", "EN")]):
        cx = mx + Inches(0.3) + i * Inches(0.95)
        fc = HONEY_BG if i == 0 else BG
        cc = HONEY_GOLD if i == 0 else SUB
        rect(slide, cx, chipset_y, Inches(0.88), Inches(0.4),
             fill=fc, radius=0.5)
        text(slide, cx, chipset_y, Inches(0.88), Inches(0.4),
             f"{flag} {lang}", size=9, bold=True, color=cc,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Text area
    rect(slide, mx + Inches(0.3), my + Inches(1.9), mw - Inches(0.6),
         Inches(1.3), fill=BG, radius=0.08)
    text(slide, mx + Inches(0.45), my + Inches(2.0), mw - Inches(0.8),
         Inches(0.3), "내용을 입력하세요...", size=11, color=SUB)
    # Buttons row
    br_y = my + Inches(3.35)
    actions = [("📷 사진", HONEY_BG, HONEY_GOLD),
                ("🎙 녹음", HONEY_BG, HONEY_GOLD),
                ("🖼 활동지", HONEY_BG, HONEY_GOLD)]
    for i, (lbl, fc, cc) in enumerate(actions):
        cx = mx + Inches(0.3) + i * Inches(1.35)
        rect(slide, cx, br_y, Inches(1.25), Inches(0.45),
             fill=fc, radius=0.3)
        text(slide, cx, br_y, Inches(1.25), Inches(0.45),
             lbl, size=10, bold=True, color=cc,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Submit
    rect(slide, mx + Inches(0.3), my + Inches(3.9), mw - Inches(0.6),
         Inches(0.45), fill=HONEY_GOLD, radius=0.3)
    text(slide, mx + Inches(0.3), my + Inches(3.9), mw - Inches(0.6),
         Inches(0.45), "카드 올리기", size=12, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def mock_tree(slide, x, y, w, h):
    # Header strip
    rect(slide, x, y, w, Inches(0.5), fill=HONEY_BG)
    text(slide, x + Inches(0.2), y + Inches(0.1), Inches(8), Inches(0.3),
         "🌳 의견 나무 · 오늘 가장 기억에 남는 것은?",
         size=12, bold=True, color=HONEY_GOLD)
    text(slide, x + w - Inches(0.6), y + Inches(0.1), Inches(0.3),
         Inches(0.3), "✕", size=14, bold=True, color=INK,
         align=PP_ALIGN.CENTER)
    # Trunk
    trunk = rect(slide, x + w // 2 - Inches(0.15), y + h - Inches(1.3),
                 Inches(0.3), Inches(1.3),
                 fill=RGBColor(0x92, 0x54, 0x0F))
    # Canopy ring of fruits
    cx = x + w // 2
    cy = y + Inches(2.7)
    import math
    colors = [HONEY_GOLD, RED, HONEY_LIGHT, GREEN, PURPLE,
              RGBColor(0xEC, 0x48, 0x99), RGBColor(0x06, 0xB6, 0xD4),
              RGBColor(0xF5, 0x9E, 0x0B)]
    for i in range(8):
        a = 2 * math.pi * i / 8
        rx = int(math.cos(a) * Inches(3.6).emu * 0.85)
        ry = int(math.sin(a) * Inches(1.6).emu * 0.9)
        fx = cx + rx - Inches(0.3)
        fy = cy + ry - Inches(0.3)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, fx, fy,
                                      Inches(0.6), Inches(0.6))
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors[i % len(colors)]
        dot.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        dot.line.width = Pt(2)
        dot.shadow.inherit = False
    # Center bee
    text(slide, cx - Inches(0.4), cy - Inches(0.4),
         Inches(0.8), Inches(0.8), "🐝",
         size=36, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def mock_game_hub(slide, x, y, w, h):
    # Header
    rect(slide, x, y, w, Inches(0.5), fill=HONEY_GOLD)
    text(slide, x + Inches(0.2), y + Inches(0.1), Inches(6), Inches(0.3),
         "🎮 소통의 방 · 게임을 골라요",
         size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    # Grid of 10 game tiles
    tiles = [
        ("🗺", "나라 맞히기"), ("🧠", "단어 기억"),
        ("🔍", "틀린 그림"), ("🧩", "문화 퍼즐"),
        ("👋", "인사말"), ("🎨", "그림 맞히기"),
        ("🔢", "숫자 탭"), ("😊", "감정 퀴즈"),
        ("🛒", "시장 역할극"), ("🍯", "단어 타워"),
    ]
    cols = 5
    rows = 2
    gx = x + Inches(0.3)
    gy = y + Inches(0.75)
    tw = (w - Inches(0.6) - Inches(0.15) * (cols - 1)) // cols
    th = (h - Inches(1.2) - Inches(0.15) * (rows - 1)) // rows
    for i, (emoji, name) in enumerate(tiles):
        r, c = divmod(i, cols)
        tx = gx + c * (tw + Inches(0.15))
        ty = gy + r * (th + Inches(0.15))
        rect(slide, tx, ty, tw, th, fill=CARD, line=HONEY_GOLD, radius=0.08)
        text(slide, tx, ty + Inches(0.1), tw, Inches(0.5),
             emoji, size=22, align=PP_ALIGN.CENTER)
        text(slide, tx, ty + th - Inches(0.4), tw, Inches(0.3),
             name, size=10, bold=True, align=PP_ALIGN.CENTER,
             color=INK)


def mock_worksheet(slide, x, y, w, h):
    # Upload card
    cw = Inches(5)
    ch = Inches(3.5)
    cx = x + (w - cw) // 2
    cy = y + Inches(0.5)
    rect(slide, cx, cy, cw, ch, fill=CARD, line=HONEY_GOLD, radius=0.05)
    text(slide, cx, cy + Inches(0.4), cw, Inches(0.5),
         "🖼 학급 활동지 사진 올리기",
         size=15, bold=True, align=PP_ALIGN.CENTER)
    rect(slide, cx + Inches(0.5), cy + Inches(1.1), cw - Inches(1),
         Inches(1.5), fill=BG, line=HONEY_GOLD, radius=0.05)
    text(slide, cx + Inches(0.5), cy + Inches(1.5), cw - Inches(1),
         Inches(0.4), "📷  사진을 여기에 놓으세요",
         size=13, color=HONEY_GOLD, bold=True, align=PP_ALIGN.CENTER)
    text(slide, cx + Inches(0.5), cy + Inches(2.0), cw - Inches(1),
         Inches(0.3), "JPG · PNG · HEIC 지원",
         size=10, color=SUB, align=PP_ALIGN.CENTER)
    rect(slide, cx + Inches(0.5), cy + Inches(2.8), cw - Inches(1),
         Inches(0.55), fill=HONEY_GOLD, radius=0.3)
    text(slide, cx + Inches(0.5), cy + Inches(2.8), cw - Inches(1),
         Inches(0.55), "✨ AI로 글자 추출하기",
         size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def mock_export_menu(slide, x, y, w, h):
    # Board background
    mock_board(slide, x, y, w, h)
    # Dropdown near the export button (top-right)
    dx = x + w - Inches(3.2)
    dy = y + Inches(0.55)
    rect(slide, dx, dy, Inches(3.0), Inches(2.2),
         fill=CARD, line=LIGHT_BORDER, radius=0.05)
    items = [("🖼  보드 전체 이미지", HONEY_GOLD),
              ("📄  PDF로 저장", INK),
              ("🎨  PPT 슬라이드", INK),
              ("📊  엑셀 CSV", INK)]
    for i, (lbl, cc) in enumerate(items):
        text(slide, dx + Inches(0.2), dy + Inches(0.15 + i * 0.5),
             Inches(2.7), Inches(0.4), lbl, size=11, bold=(i == 0),
             color=cc, anchor=MSO_ANCHOR.MIDDLE)


# ---------- Slide builders ----------
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, HONEY_BG)
    # Big center bee
    text(s, 0, Inches(1.2), SW, Inches(1.6), "🐝",
         size=120, align=PP_ALIGN.CENTER)
    text(s, 0, Inches(3.0), SW, Inches(0.9), "꿀벌 소통창",
         size=54, bold=True, color=INK, align=PP_ALIGN.CENTER)
    text(s, 0, Inches(4.0), SW, Inches(0.5),
         "선생님용 사용 매뉴얼",
         size=22, color=HONEY_GOLD, bold=True, align=PP_ALIGN.CENTER)
    text(s, 0, Inches(5.0), SW, Inches(0.4),
         "다문화 학급을 위한 실시간 다국어 소통 보드",
         size=13, color=SUB, align=PP_ALIGN.CENTER)
    # Bottom info
    rect(s, 0, SH - Inches(0.8), SW, Inches(0.8), fill=INK)
    text(s, 0, SH - Inches(0.8), SW, Inches(0.8),
         "기능별 따라하기 · 15개 언어 자동 번역 · 음성 읽기 · AI 글자 인식",
         size=12, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_toc(total_pages):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    text(s, Inches(0.8), Inches(0.6), Inches(12), Inches(0.7),
         "📑  목차",
         size=32, bold=True, color=INK)
    text(s, Inches(0.8), Inches(1.3), Inches(12), Inches(0.4),
         "아래 8개 장으로 구성되어 있어요. 필요한 부분만 골라 보셔도 좋아요.",
         size=13, color=SUB)
    chapters = [
        ("01", "시작하기",     "방 만들기 · 학생 초대 · 접속하기"),
        ("02", "보드 둘러보기", "카드 · 색상 · 실시간 업데이트"),
        ("03", "카드 작성",     "글 · 사진 · 음성 녹음"),
        ("04", "번역과 음성",   "15개 언어 자동 번역 · TTS 듣기"),
        ("05", "의견 나누기",   "🌳 나무 소통방 · 비공개 의견 모으기"),
        ("06", "학급 활동지",   "📷 사진 한 장으로 AI 글자 추출"),
        ("07", "소통의 방",     "🎮 10가지 짝 맞추기·퀴즈 게임"),
        ("08", "내보내기",     "🖼 보드 이미지 · 📄 PDF · 📊 CSV"),
    ]
    cols = 2
    gx = Inches(0.8)
    gy = Inches(2.0)
    cw = Inches(5.9)
    ch = Inches(1.15)
    gap = Inches(0.2)
    for i, (no, title, desc) in enumerate(chapters):
        r, c = divmod(i, cols)
        x = gx + c * (cw + gap)
        y = gy + r * (ch + gap)
        rect(s, x, y, cw, ch, fill=CARD, line=LIGHT_BORDER, radius=0.04)
        # Number badge
        rect(s, x + Inches(0.2), y + Inches(0.22), Inches(0.7),
             Inches(0.7), fill=HONEY_BG, radius=0.15)
        text(s, x + Inches(0.2), y + Inches(0.22), Inches(0.7),
             Inches(0.7), no, size=18, bold=True, color=HONEY_GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + Inches(1.05), y + Inches(0.2), cw - Inches(1.2),
             Inches(0.4), title, size=15, bold=True, color=INK)
        text(s, x + Inches(1.05), y + Inches(0.6), cw - Inches(1.2),
             Inches(0.4), desc, size=11, color=SUB)
    footer(s, 2, total_pages)


def slide_intro(total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    page_header(s, "00", "이 앱은 어떤 앱인가요?",
                "다문화 학급에서 언어의 벽 없이 소통할 수 있게 해주는 보드 앱입니다.")
    # 3-column feature cards
    items = [
        ("🌍", "15개 언어", "한국어·영어·중국어·베트남어·일본어 외 10개 언어 자동 번역"),
        ("🔊", "음성으로 듣기", "읽지 못하는 글도 소리로 들을 수 있어 모든 학생이 참여 가능"),
        ("🎨", "쉬운 작성", "텍스트·사진·녹음·활동지 촬영 — 스마트폰만으로 OK"),
        ("🌳", "의견 모으기", "선생님이 주제를 내면 학생들이 비공개로 답을 올리고 나무 열매로 공개"),
        ("🎮", "재미있는 게임", "10가지 미니 게임으로 다른 언어를 놀이처럼 익혀요"),
        ("🖼", "한 장으로 저장", "하루 활동을 이미지·PDF·CSV로 내려받아 기록 보관"),
    ]
    gx = Inches(0.7)
    gy = Inches(1.7)
    cw = Inches(4.0)
    ch = Inches(2.55)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)
    cols = 3
    for i, (emoji, title, desc) in enumerate(items):
        r, c = divmod(i, cols)
        x = gx + c * (cw + gap_x)
        y = gy + r * (ch + gap_y)
        rect(s, x, y, cw, ch, fill=CARD, line=LIGHT_BORDER, radius=0.05)
        text(s, x + Inches(0.3), y + Inches(0.2), Inches(1), Inches(0.6),
             emoji, size=30)
        text(s, x + Inches(0.3), y + Inches(1.0), cw - Inches(0.6),
             Inches(0.4), title, size=16, bold=True, color=INK)
        text(s, x + Inches(0.3), y + Inches(1.5), cw - Inches(0.6),
             Inches(1.0), desc, size=11, color=SUB)
    footer(s, 3, total)


def slide_chapter_cover(no, title, emoji, desc, total, page):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, HONEY_BG)
    text(s, 0, Inches(1.8), SW, Inches(1.2), emoji,
         size=110, align=PP_ALIGN.CENTER)
    text(s, 0, Inches(3.4), SW, Inches(0.5),
         f"Chapter {no}",
         size=18, bold=True, color=HONEY_GOLD, align=PP_ALIGN.CENTER)
    text(s, 0, Inches(3.9), SW, Inches(0.8), title,
         size=40, bold=True, color=INK, align=PP_ALIGN.CENTER)
    text(s, 0, Inches(4.9), SW, Inches(0.5), desc,
         size=14, color=SUB, align=PP_ALIGN.CENTER)
    footer(s, page, total)


def two_col_slide(no, chapter_title, subtitle, mock_fn, steps,
                  total, page, tips=None):
    """Left: screen mock. Right: numbered steps."""
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    page_header(s, no, chapter_title, subtitle)
    # Left mock
    mx = Inches(0.6)
    my = Inches(1.6)
    mw = Inches(6.6)
    mh = Inches(4.5)
    mock_screen(s, mx, my, mw, mh, mock_fn)
    # Right steps
    sx = Inches(7.5)
    sy = Inches(1.6)
    sw = Inches(5.2)
    text(s, sx, sy, sw, Inches(0.5), "✅ 이렇게 하세요",
         size=16, bold=True, color=HONEY_GOLD)
    for i, (t, b) in enumerate(steps):
        numbered_step(s, sx, sy + Inches(0.6 + i * 1.05), sw,
                      i + 1, t, b)
    # Tips box
    if tips:
        ty = Inches(6.3)
        rect(s, mx, ty, SW - Inches(1.2), Inches(0.85),
             fill=HONEY_BG, radius=0.1)
        text(s, mx + Inches(0.2), ty + Inches(0.08),
             SW - Inches(1.6), Inches(0.3),
             "💡  TIP", size=11, bold=True, color=HONEY_GOLD)
        text(s, mx + Inches(0.2), ty + Inches(0.4),
             SW - Inches(1.6), Inches(0.4),
             tips, size=12, color=INK)
    footer(s, page, total)


def slide_language_list(total, page):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    page_header(s, "04", "지원하는 15개 언어",
                "학생이 자기 모국어로 쓰면 다른 학생 화면에는 자동으로 번역되어 보입니다.")
    langs = [
        ("🇰🇷", "한국어"), ("🇻🇳", "Tiếng Việt"),
        ("🇨🇳", "中文"),  ("🇺🇸", "English"),
        ("🇯🇵", "日本語"), ("🇵🇭", "Filipino"),
        ("🇹🇭", "ภาษาไทย"), ("🇲🇳", "Монгол"),
        ("🇺🇿", "O'zbek"), ("🇷🇺", "Русский"),
        ("🇮🇩", "Bahasa Indonesia"), ("🇰🇭", "ខ្មែរ"),
        ("🇲🇲", "မြန်မာ"), ("🇳🇵", "नेपाली"),
        ("🇮🇳", "हिन्दी"),
    ]
    cols = 5
    rows = 3
    gx = Inches(0.7)
    gy = Inches(1.7)
    tw = Inches(2.4)
    th = Inches(1.45)
    gap = Inches(0.15)
    for i, (flag, name) in enumerate(langs):
        r, c = divmod(i, cols)
        x = gx + c * (tw + gap)
        y = gy + r * (th + gap)
        rect(s, x, y, tw, th, fill=CARD, line=LIGHT_BORDER, radius=0.06)
        text(s, x, y + Inches(0.15), tw, Inches(0.6), flag,
             size=30, align=PP_ALIGN.CENTER)
        text(s, x, y + Inches(0.9), tw, Inches(0.4), name,
             size=13, bold=True, align=PP_ALIGN.CENTER, color=INK)
    # Bottom emphasis
    by = Inches(6.2)
    rect(s, Inches(0.7), by, SW - Inches(1.4), Inches(0.9),
         fill=HONEY_BG, radius=0.1)
    text(s, 0, by + Inches(0.1), SW, Inches(0.3),
         "한 번 올린 카드를 학생마다 다른 언어로 자동 표시 · 읽기 어려운 글자는 🔊 버튼으로 음성으로 들을 수 있어요",
         size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
    footer(s, page, total)


def slide_games_grid(total, page):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    page_header(s, "07", "소통의 방 · 10가지 게임",
                "게임을 시작할 때 2개 언어를 고르면, 두 언어로 같이 나옵니다.")
    games = [
        ("🗺", "나라 맞히기", "국기를 보고 나라 이름 고르기"),
        ("🧠", "단어 기억",    "같은 뜻의 카드 짝 맞추기"),
        ("🔍", "틀린 그림 찾기","두 그림에서 다른 곳 찾기"),
        ("🧩", "문화 퍼즐",    "한복·쌀국수·게르 조각 맞추기"),
        ("👋", "인사말 릴레이","들은 인사말 같은 뜻 고르기"),
        ("🎨", "그림 맞히기",  "꿀벌 그림 보고 단어 맞히기"),
        ("🔢", "숫자 탭",      "30초 안에 숫자 빨리 누르기"),
        ("😊", "감정 퀴즈",    "상황에 맞는 이모지 고르기"),
        ("🛒", "시장 역할극",  "대화 상황에 맞게 답하기"),
        ("🍯", "단어 타워",    "맞히면 꿀단지 쌓기 / 3번 틀리면 끝"),
    ]
    cols = 5
    rows = 2
    gx = Inches(0.6)
    gy = Inches(1.7)
    tw = Inches(2.4)
    th = Inches(2.5)
    gap = Inches(0.15)
    for i, (emoji, name, desc) in enumerate(games):
        r, c = divmod(i, cols)
        x = gx + c * (tw + gap)
        y = gy + r * (th + gap)
        rect(s, x, y, tw, th, fill=CARD, line=HONEY_GOLD, radius=0.08)
        text(s, x, y + Inches(0.25), tw, Inches(0.8), emoji,
             size=34, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.15), y + Inches(1.25), tw - Inches(0.3),
             Inches(0.4), name, size=13, bold=True,
             align=PP_ALIGN.CENTER, color=INK)
        text(s, x + Inches(0.15), y + Inches(1.7), tw - Inches(0.3),
             Inches(0.7), desc, size=10, color=SUB,
             align=PP_ALIGN.CENTER)
    footer(s, page, total)


def slide_qa(total, page):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    page_header(s, "✅", "자주 묻는 질문 (FAQ)",
                "문제가 생겼을 때 먼저 확인해 보세요.")
    qas = [
        ("학생이 방에 접속하지 못해요",
         "방 코드를 다시 확인하세요. 대문자·소문자는 구분하지 않아요. 인터넷 연결 상태도 확인해주세요."),
        ("번역이 안 나와요",
         "카드를 만들 때 본인의 '원래 언어'를 국기 칩에서 한 번 골라야 다른 학생 화면에서 자동 번역돼요."),
        ("음성이 안 들려요",
         "브라우저에서 소리 권한이 차단된 경우가 있어요. 주소창 왼쪽 자물쇠 아이콘을 눌러 소리를 '허용'으로 바꿔주세요."),
        ("활동지 사진에서 글자가 잘못 나와요",
         "글씨가 흐리거나 빛 반사가 있을 때 정확도가 떨어집니다. 정면에서 밝게 찍으면 훨씬 좋아요."),
        ("학생이 카드를 지울 수 있나요?",
         "본인이 쓴 카드는 본인이 지울 수 있어요. 선생님은 모든 카드를 관리할 수 있습니다."),
        ("수업이 끝난 뒤 기록을 남기고 싶어요",
         "보드 오른쪽 위 '내보내기' → '🖼 보드 전체 이미지'를 누르면 전체가 한 장으로 저장됩니다."),
    ]
    gx = Inches(0.6)
    gy = Inches(1.6)
    cw = Inches(6.0)
    ch = Inches(1.7)
    gap_x = Inches(0.15)
    gap_y = Inches(0.15)
    for i, (q, a) in enumerate(qas):
        r, c = divmod(i, 2)
        x = gx + c * (cw + gap_x)
        y = gy + r * (ch + gap_y)
        rect(s, x, y, cw, ch, fill=CARD, line=LIGHT_BORDER, radius=0.04)
        rect(s, x + Inches(0.2), y + Inches(0.2), Inches(0.45),
             Inches(0.45), fill=HONEY_GOLD, radius=0.5)
        text(s, x + Inches(0.2), y + Inches(0.2), Inches(0.45),
             Inches(0.45), "Q", size=13, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + Inches(0.8), y + Inches(0.15), cw - Inches(1),
             Inches(0.5), q, size=13, bold=True, color=INK)
        text(s, x + Inches(0.8), y + Inches(0.6), cw - Inches(1),
             Inches(1.0), a, size=11, color=SUB)
    footer(s, page, total)


def slide_end(total, page):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, HONEY_GOLD)
    text(s, 0, Inches(2.2), SW, Inches(1.3), "🐝",
         size=110, align=PP_ALIGN.CENTER)
    text(s, 0, Inches(3.7), SW, Inches(0.8),
         "수고하셨어요!",
         size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER)
    text(s, 0, Inches(4.6), SW, Inches(0.5),
         "오늘도 우리 교실을 따뜻한 꿀벌 소통창으로 채워 주세요.",
         size=16, color=HONEY_BG, align=PP_ALIGN.CENTER)
    footer(s, page, total)


# ---------- Build ----------
# We know total pages ahead of time to number footers correctly.
TOTAL = 19

# 1. Cover
slide_cover()
# 2. TOC
slide_toc(TOTAL)
# 3. Intro
slide_intro(TOTAL)

# Chapter 01 - 시작하기
slide_chapter_cover("01", "시작하기", "🚪",
                     "방을 만들고 학생을 초대하는 첫 단계", TOTAL, 4)

two_col_slide("01", "새 방 만들기",
    "선생님이 먼저 방을 만들고, 학생에게 방 코드를 공유합니다.",
    mock_cover_landing,
    steps=[
        ("앱 접속하기",
         "인터넷 주소(URL)로 접속하면 꿀벌 마스코트가 있는 첫 화면이 뜹니다."),
        ("‘새 방 만들기’ 버튼",
         "꿀색 버튼을 누르면 오늘의 주제·방 이름을 입력하는 화면이 나옵니다."),
        ("주제 입력",
         "예: ‘우리 반 자기소개’ · ‘오늘 점심 후기’처럼 간단한 문장이 좋습니다."),
        ("방 코드 확인",
         "화면에 6자리 코드가 뜹니다. 이것을 학생들에게 알려주세요."),
    ],
    tips="첫 화면에서는 로그인 없이 바로 시작할 수 있어요. 방 코드만 있으면 누구든 들어올 수 있습니다.",
    total=TOTAL, page=5)

two_col_slide("01", "학생 초대하기",
    "학생은 휴대폰으로 같은 사이트에 접속해 방 코드를 입력합니다.",
    mock_cover_landing,
    steps=[
        ("사이트 주소 공유",
         "주소를 QR코드로 만들거나 칠판에 적어 주세요."),
        ("‘방 코드 입력’ 선택",
         "학생은 첫 화면에서 하얀색 ‘방 코드 입력’을 누릅니다."),
        ("이름·국기 고르기",
         "학생이 자기 이름과 모국어 국기를 한 번 고릅니다."),
        ("보드 입장 완료",
         "완료하면 바로 보드가 열립니다. 카드를 써서 올릴 수 있어요."),
    ],
    tips="방 코드는 대문자·소문자를 구분하지 않아 편합니다. 공백만 주의해 주세요.",
    total=TOTAL, page=6)

# Chapter 02 - 보드 둘러보기
slide_chapter_cover("02", "보드 둘러보기", "📌",
                     "카드가 실시간으로 쌓이는 소통 공간", TOTAL, 7)

two_col_slide("02", "보드 화면 구성",
    "위쪽에 방 이름과 주요 버튼, 아래쪽에 학생들이 올린 카드들이 보입니다.",
    mock_board,
    steps=[
        ("보드 제목",
         "방 이름과 코드가 왼쪽 위에 표시돼요. 수업 중 늘 보여서 편리해요."),
        ("카드 추가 버튼",
         "보라색 ‘+ 카드’ 버튼으로 누구나 새 카드를 올릴 수 있어요."),
        ("게임·내보내기",
         "오른쪽 위 ‘🎮 게임’과 ‘⬇ 내보내기’로 학급 활동을 더 풍부하게."),
        ("실시간 업데이트",
         "새 카드는 모든 학생·선생님 화면에 곧바로 나타납니다."),
    ],
    tips="카드 색깔은 모국어 국기에 따라 조금씩 달라져 누가 올렸는지 직관적으로 구분할 수 있어요.",
    total=TOTAL, page=8)

# Chapter 03 - 카드 작성
slide_chapter_cover("03", "카드 작성", "✏️",
                     "글·사진·음성을 자유롭게 담는 방법", TOTAL, 9)

two_col_slide("03", "카드 쓰기 창",
    "‘+ 카드’를 누르면 작성 창이 뜹니다. 언어·내용·첨부를 순서대로 채우세요.",
    mock_card_modal,
    steps=[
        ("이름 입력",
         "학생 이름이 자동으로 들어가 있어요. 다르게 나오길 원하면 바꿔 주세요."),
        ("모국어 국기 선택",
         "처음 한 번 고르면 그 후에는 자동 기억됩니다."),
        ("내용 입력",
         "평소 쓰는 언어로 그냥 쓰면, 다른 학생에게는 번역되어 보입니다."),
        ("사진·녹음·활동지",
         "아래 세 버튼으로 원하는 자료를 첨부할 수 있어요."),
        ("‘카드 올리기’ 누르기",
         "꿀색 버튼을 누르면 보드에 바로 올라갑니다."),
    ],
    tips="사진은 자동으로 압축되어 빠르게 올라갑니다. 녹음은 60초까지 지원해요.",
    total=TOTAL, page=10)

# Chapter 04 - 번역과 음성
slide_chapter_cover("04", "번역과 음성", "🌐",
                     "15개 언어를 자유롭게 오가는 소통", TOTAL, 11)
slide_language_list(TOTAL, 12)

# Chapter 05 - 의견 나누기
slide_chapter_cover("05", "의견 나누기 — 나무 소통방", "🌳",
                     "선생님이 질문하면 학생이 비공개로 답하고, 함께 열매로 공개", TOTAL, 13)
two_col_slide("05", "나무 소통방 만들기",
    "선생님만 쓸 수 있는 기능입니다. ‘소통방’ 버튼으로 질문을 던져 보세요.",
    mock_tree,
    steps=[
        ("주제 입력",
         "예: ‘오늘 가장 기억에 남는 것은?’ 같은 열린 질문이 좋아요."),
        ("학생이 비공개로 답",
         "학생 답은 잠금 아이콘이 붙어 있어 다른 학생에게 보이지 않아요."),
        ("‘다 같이 열기’",
         "선생님이 열면 학생들의 답이 열매가 되어 나무에 한꺼번에 열립니다."),
        ("열매 클릭",
         "어느 열매든 눌러 해당 의견을 모두에게 보여줄 수 있습니다."),
    ],
    tips="감정·소감·한마디처럼 짧은 대답이 어울려요. 긴 토론은 일반 카드를 쓰시길 권해요.",
    total=TOTAL, page=14)

# Chapter 06 - 활동지
slide_chapter_cover("06", "학급 활동지 촬영", "📷",
                     "사진 한 장으로 손글씨·프린트를 글자로 바꿔요", TOTAL, 15)
two_col_slide("06", "활동지 올리기",
    "카드 작성 창에서 ‘🖼 활동지’를 누르면 AI가 사진 속 글자를 자동으로 읽어 줍니다.",
    mock_worksheet,
    steps=[
        ("활동지를 찍거나 선택",
         "인쇄물·공책·판서 사진 어떤 것도 가능합니다."),
        ("AI가 글자 인식",
         "사진을 올린 뒤 ‘✨ AI로 글자 추출하기’를 누르면 잠시 후 텍스트가 나타납니다."),
        ("필요한 부분만 수정",
         "인식이 애매한 부분은 손으로 고쳐도 돼요."),
        ("카드로 올리기",
         "그대로 카드에 담으면 다른 언어로도 번역되어 모두가 볼 수 있어요."),
    ],
    tips="글씨는 정면에서 밝게 찍으면 인식률이 훨씬 좋아집니다. 너무 작은 글씨는 살짝 확대해 주세요.",
    total=TOTAL, page=16)

# Chapter 07 - 소통의 방
slide_chapter_cover("07", "소통의 방 — 10가지 게임", "🎮",
                     "언어를 놀이처럼 — 2개 언어를 고르면 함께 표시됩니다", TOTAL, 17)
slide_games_grid(TOTAL, 18)

# Chapter 08 - 내보내기
slide_chapter_cover("08", "내보내기", "⬇",
                     "오늘 활동을 한 장의 이미지·파일로 남겨요", TOTAL, 19)

# Actually we said TOTAL=19 but let's adjust. Let me recount.
# Added so far:
# 1 cover, 2 toc, 3 intro,
# 4 ch1 cover, 5 ch1 p1, 6 ch1 p2,
# 7 ch2 cover, 8 ch2 p1,
# 9 ch3 cover, 10 ch3 p1,
# 11 ch4 cover, 12 ch4 p1,
# 13 ch5 cover, 14 ch5 p1,
# 15 ch6 cover, 16 ch6 p1,
# 17 ch7 cover, 18 ch7 p1,
# 19 ch8 cover
# We still need ch8 page + FAQ + end.
# Let me extend TOTAL.
# Going to rebuild after counting properly.
# Save file and re-run — we'll recompute.

two_col_slide("08", "보드 내보내기",
    "오른쪽 위 ‘⬇ 내보내기’에서 원하는 형식을 고르면 즉시 다운로드됩니다.",
    mock_export_menu,
    steps=[
        ("‘⬇ 내보내기’ 누르기",
         "오른쪽 위 하얀 버튼을 누르면 작은 메뉴가 내려와요."),
        ("형식 고르기",
         "이미지·PDF·PPT·CSV 중 하나를 선택합니다."),
        ("‘🖼 보드 전체 이미지’ 추천",
         "학급 게시용으로 가장 깔끔해요. 보드 전체가 한 장으로 저장됩니다."),
        ("저장 위치 확인",
         "브라우저의 ‘다운로드’ 폴더에 파일이 저장됩니다."),
    ],
    tips="PDF는 인쇄용, CSV는 엑셀에서 분석할 때, 이미지는 교실 게시판·가정통신문용으로 적합해요.",
    total=TOTAL, page=20)

# FAQ
slide_qa(TOTAL, 21)

# End
slide_end(TOTAL, 22)


# ---- Patch footers with the right total ----
# We set TOTAL=19 initially but actual total = 22. The footer shows wrong denominator.
# Fix by overwriting all footer text runs.
ACTUAL_TOTAL = len(prs.slides)
for idx, slide in enumerate(prs.slides, start=1):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.text and "/ " in r.text and r.text.strip().endswith("19"):
                    parts = r.text.split("/")
                    r.text = f"{parts[0]}/ {ACTUAL_TOTAL}"

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {ACTUAL_TOTAL}")
